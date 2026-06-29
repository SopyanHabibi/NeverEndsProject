import os
import re
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

UKURAN_CHUNK_KATA = 300  # ~300 kata per chunk, cukup kecil buat context 4096 token

def ekstrak_teks(filepath: str, ekstensi: str) -> str:
    """Ekstrak teks mentah dari PDF/DOCX/PPTX."""
    ekstensi = ekstensi.lower()
    try:
        if ekstensi == "pdf":
            reader = PdfReader(filepath)
            return "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ekstensi == "docx":
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ekstensi == "pptx":
            prs = Presentation(filepath)
            teks_slide = []
            for i, slide in enumerate(prs.slides, 1):
                isi_slide = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        isi_slide.append(shape.text)
                teks_slide.append("\n".join(isi_slide))
            return "\n\n".join(teks_slide)

        else:
            return ""
    except Exception as e:
        print(f"[DOKUMEN ERROR] Gagal ekstrak {filepath}: {e}")
        return ""

def pecah_jadi_chunks(teks: str) -> list:
    """Pecah teks panjang jadi potongan ~300 kata, biar gampang dicari relevansinya nanti."""
    kata = teks.split()
    chunks = []
    for i in range(0, len(kata), UKURAN_CHUNK_KATA):
        potongan = " ".join(kata[i:i + UKURAN_CHUNK_KATA])
        if potongan.strip():
            chunks.append(potongan)
    return chunks

def _skor_overlap(pertanyaan: str, chunk: str) -> int:
    """Skor sederhana: hitung berapa kata di pertanyaan yang juga muncul di chunk ini."""
    kata_pertanyaan = set(re.findall(r'\w+', pertanyaan.lower()))
    kata_chunk = set(re.findall(r'\w+', chunk.lower()))
    # Buang kata umum yang gak informatif biar gak nge-bias hasil
    stopwords = {"the", "a", "an", "is", "are", "what", "how", "yang", "dan", "di", "ini", "itu", "apa", "untuk"}
    kata_pertanyaan -= stopwords
    return len(kata_pertanyaan & kata_chunk)

def pilih_chunks_relevan(chunks: list, pertanyaan: str, top_n: int = 3) -> str:
    """Cari N chunk paling relevan berdasarkan overlap kata dengan pertanyaan."""
    if not chunks:
        return ""
    skor_list = [(c, _skor_overlap(pertanyaan, c["konten"])) for c in chunks]
    skor_list.sort(key=lambda x: x[1], reverse=True)
    terpilih = [c for c, skor in skor_list[:top_n] if skor > 0]

    if not terpilih:
        # Gak ada yang match sama sekali -> kasih chunk pertama aja sebagai fallback
        terpilih = chunks[:1]

    terpilih.sort(key=lambda c: c["index"])  # urutkan balik sesuai posisi asli di dokumen
    return "\n\n".join(c["konten"] for c in terpilih)

def pilih_chunks_sample(chunks: list, max_chunks: int = 5) -> str:
    """Buat summarize/quiz: ambil sample merata dari seluruh dokumen (awal-tengah-akhir),
    bukan cuma 1 bagian, biar representasinya menyeluruh."""
    if not chunks:
        return ""
    total = len(chunks)
    if total <= max_chunks:
        terpilih = chunks
    else:
        step = total / max_chunks
        indeks_terpilih = [int(i * step) for i in range(max_chunks)]
        terpilih = [chunks[i] for i in indeks_terpilih]
    return "\n\n".join(c["konten"] for c in terpilih)